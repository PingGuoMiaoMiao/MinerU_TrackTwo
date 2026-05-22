from pathlib import Path

import httpx
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from src.config import settings
from src.tools.mineru import MinerUError, mineru_client


def parse_text(text: str) -> tuple[str, dict]:
    return text, {"format": "text", "chars": len(text)}


def parse_file(path: Path) -> tuple[str, dict]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix in {".docx", ".doc"}:
        return parse_docx(path)
    if suffix in {".pptx", ".ppt"}:
        return parse_pptx(path)
    if suffix in {".html", ".htm"}:
        return parse_html(path.read_text(encoding="utf-8", errors="ignore"), source=str(path))
    return path.read_text(encoding="utf-8", errors="ignore"), {
        "format": suffix.lstrip(".") or "text",
        "filename": path.name,
        "bytes": path.stat().st_size,
    }


async def parse_file_auto(path: Path, data_id: str | None = None) -> tuple[str, dict]:
    suffix = path.suffix.lower()
    if settings.mineru_mode.lower() != "off" and mineru_client.supports_path(path) and suffix not in {".txt", ".md", ".html", ".htm"}:
        try:
            return await mineru_client.parse_file(path, data_id=data_id)
        except MinerUError as exc:
            text, metadata = parse_file(path)
            metadata["parser"] = "local_fallback"
            metadata["mineru_error"] = str(exc)
            return text, metadata
    text, metadata = parse_file(path)
    metadata.setdefault("parser", "local")
    return text, metadata


def parse_pdf(path: Path) -> tuple[str, dict]:
    reader = PdfReader(str(path))
    pages = []
    for page_no, page in enumerate(reader.pages, start=1):
        pages.append(f"[page {page_no}]\n{page.extract_text() or ''}")
    return "\n\n".join(pages), {
        "format": "pdf",
        "filename": path.name,
        "pages": len(reader.pages),
        "bytes": path.stat().st_size,
    }


def parse_docx(path: Path) -> tuple[str, dict]:
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    table_text = []
    for table_no, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        table_text.append(f"[table {table_no}]\n" + "\n".join(rows))
    text = "\n".join(paragraphs + table_text)
    return text, {
        "format": "docx",
        "filename": path.name,
        "paragraphs": len(paragraphs),
        "tables": len(doc.tables),
        "bytes": path.stat().st_size,
    }


def parse_pptx(path: Path) -> tuple[str, dict]:
    prs = Presentation(str(path))
    slides = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append(f"[slide {slide_no}]\n" + "\n".join(parts))
    return "\n\n".join(slides), {
        "format": "pptx",
        "filename": path.name,
        "slides": len(prs.slides),
        "bytes": path.stat().st_size,
    }


def parse_html(html: str, source: str = "html") -> tuple[str, dict]:
    soup = BeautifulSoup(html, "html.parser")
    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return text, {"format": "html", "source": source, "title": title, "chars": len(text)}


async def parse_url(url: str) -> tuple[str, dict]:
    if settings.mineru_mode.lower() != "off":
        try:
            return await mineru_client.parse_url(url)
        except MinerUError:
            pass

    async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
        response = await client.get(url)
        response.raise_for_status()
    text, metadata = parse_html(response.text, source=url)
    metadata["status_code"] = response.status_code
    metadata["content_type"] = response.headers.get("content-type", "")
    return text, metadata
